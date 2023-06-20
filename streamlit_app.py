import streamlit as st
import pandas as pd
import numpy as np
from pandasai import PandasAI
from pandasai.llm.openai import OpenAI
from PIL import Image
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

image = Image.open('exl.png')


llm = OpenAI(api_token=st.secrets["chat_gpt_key"])
pandas_ai = PandasAI(llm, conversational=False)#, enforce_privacy = True)

df = pd.read_csv('data.csv')
    
with st.sidebar:
	st.image(image, width = 150)
	st.header('Graphs')
	st.write('Create 2*2 graph grid')


st.subheader("Data" )
st.dataframe(df.head())

with st.form("Form1"):
	option = st.selectbox("Select graph", ("Graph1","Graph2","Graph3","Graph4"))
	graph = st.text_input(label ="Graph")
	title = st.text_input(label ="Title")
	submitted1 = st.form_submit_button("Submit")
	
	
	if submitted1:
		fig, x= plt.subplots()
		response = pandas_ai(df, prompt=f"""Plot {graph} and use '{title}' as chart title""")
		st.pyplot(fig)
		fig.savefig(option+'.png')



		st.write("---")
		st.subheader("Download")
		slide_header = st.text_input(label ="Header")
		if slide_header:
			# Create a new PowerPoint presentation
			presentation = Presentation()
			
			# Define the image file paths
			image_paths = ['Graph1.png', 'Graph2.png', 'Graph3.png', 'Graph4.png']
			
			# Create a slide with a 2x2 image grid
			slide_layout = presentation.slide_layouts[6]  # Use slide layout with 2 content placeholders
			slide = presentation.slides.add_slide(slide_layout)
			shapes = slide.shapes
			# Define the dimensions of each graph
			width = Inches(4)
			height = Inches(3)
			
			text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), presentation.slide_width, Inches(0.5)).text_frame
			text_box.text = slide_header
			
			font = text_box.paragraphs[0].runs[0].font
			font.name = 'Calibri Light (Headings)'
			font.size = Pt(24)
			font.bold = True
			font.color.rgb = RGBColor(251, 77, 10)
			
			# Add the first graph to the slide
			left = Inches(0.5)
			top = Inches(0.8)
			pic = slide.shapes.add_picture('graph1.png', left, top, width=width, height=height)
			
			# Add the second graph to the slide
			left = Inches(5.5)
			pic = slide.shapes.add_picture('graph2.png', left, top, width=width, height=height)
			
			# Add the third graph to the slide
			left = Inches(0.5)
			top = Inches(3.8)
			pic = slide.shapes.add_picture('graph3.png', left, top, width=width, height=height)
			
			# Add the fourth graph to the slide
			left = Inches(5.5)
			pic = slide.shapes.add_picture('graph4.png', left, top, width=width, height=height)
			
			# Add page number and trademark sign
			slide_number = presentation.slide_height - Inches(0.5)
			slide_width = presentation.slide_width - Inches(2)
			text_box = slide.shapes.add_textbox(slide_width - Inches(2.5), slide_number, Inches(2), Inches(0.5)).text_frame
			text_box.text = "\u00A9 2023 EXL Service, Inc. All rights reserved."
			pic = slide.shapes.add_picture('exl.png', presentation.slide_width - Inches(1.5), presentation.slide_height - Inches(0.5), width=Inches(0.8), height=Inches(0.32))
			
			font = text_box.paragraphs[0].runs[0].font
			font.name = 'Calibri Light (Headings)'
			font.size = Pt(10)
			
			presentation.save('image_grid.pptx')
		    
			with open("image_grid.pptx", "rb") as file:
				st.download_button(
				    label="Download data",
				    data=file,
				    file_name='image_grid.pptx'
				    )
		st.write("---")
